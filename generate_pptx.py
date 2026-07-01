import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(overline_text, title_text, subtitle_text, col1_h, col1_p1, col1_p2, col2_h, col2_p1, col2_p2, col2_box=""):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Overline
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = overline_text
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(212, 175, 55) # Gold
    
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = title_text
    p2.font.bold = True
    p2.font.name = "Georgia"
    p2.font.size = Pt(40)
    p2.font.color.rgb = RGBColor(16, 42, 67) # Navy
    
    # Subtitle
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.add_paragraph()
    p3.text = subtitle_text
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(72, 101, 129)
    
    # Column 1 Box
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(2.8), Inches(5.6), Inches(4.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    
    # Column 1 Content
    tx_col1 = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(5.2), Inches(3.8))
    tf_col1 = tx_col1.text_frame
    tf_col1.word_wrap = True
    
    ph = tf_col1.add_paragraph()
    ph.text = col1_h
    ph.font.bold = True
    ph.font.name = "Georgia"
    ph.font.size = Pt(24)
    ph.font.color.rgb = RGBColor(16, 42, 67)
    
    p_b1 = tf_col1.add_paragraph()
    p_b1.text = "\n" + col1_p1
    p_b1.font.size = Pt(18)
    p_b1.font.color.rgb = RGBColor(51, 78, 104)
    
    p_b2 = tf_col1.add_paragraph()
    p_b2.text = "\n" + col1_p2
    p_b2.font.size = Pt(18)
    p_b2.font.color.rgb = RGBColor(51, 78, 104)
    
    # Column 2 Box
    shape2 = slide.shapes.add_shape(1, Inches(6.7), Inches(2.8), Inches(5.6), Inches(4.2))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(248, 249, 250)
    shape2.line.color.rgb = RGBColor(226, 232, 240)
    
    # Column 2 Content
    tx_col2 = slide.shapes.add_textbox(Inches(6.9), Inches(3.0), Inches(5.2), Inches(3.8))
    tf_col2 = tx_col2.text_frame
    tf_col2.word_wrap = True
    
    ph2 = tf_col2.add_paragraph()
    ph2.text = col2_h
    ph2.font.bold = True
    ph2.font.name = "Georgia"
    ph2.font.size = Pt(24)
    ph2.font.color.rgb = RGBColor(16, 42, 67)
    
    p_c1 = tf_col2.add_paragraph()
    p_c1.text = "\n" + col2_p1
    p_c1.font.size = Pt(18)
    p_c1.font.color.rgb = RGBColor(51, 78, 104)
    
    p_c2 = tf_col2.add_paragraph()
    p_c2.text = "\n" + col2_p2
    p_c2.font.size = Pt(18)
    p_c2.font.color.rgb = RGBColor(51, 78, 104)
    
    if col2_box:
        # Highlight Box inside Col 2
        box3 = slide.shapes.add_shape(1, Inches(6.9), Inches(5.8), Inches(5.2), Inches(1.0))
        box3.fill.solid()
        box3.fill.fore_color.rgb = RGBColor(16, 42, 67)
        tx_box3 = slide.shapes.add_textbox(Inches(7.1), Inches(6.0), Inches(4.8), Inches(0.8))
        tf_box3 = tx_box3.text_frame
        tf_box3.word_wrap = True
        pb3 = tf_box3.add_paragraph()
        pb3.text = col2_box
        pb3.font.italic = True
        pb3.font.size = Pt(16)
        pb3.font.color.rgb = RGBColor(255, 255, 255)

# Slide 1
add_slide(
    "E N G I N E E R I N G   D E P T H   ·   P I L L A R   2",
    "Sub-second velocity over 100k events.",
    "Why we offloaded incremental stream processing to Rust.",
    "The Python Bottleneck",
    "Catching \"velocity bursts\" (e.g., 50 micro-deposits in 5 minutes) requires sliding-window analytics over live streams.",
    "Doing this in pure Python over a live Kafka feed is CPU-heavy and creates a massive bottleneck under true PSB transaction loads.",
    "The Pathway (Rust) Engine",
    "We offloaded this specific computation to Pathway, a Rust-backed incremental streaming engine.",
    "This allows us to keep our core logic in Python while leveraging Rust's bare-metal speed for the heavy incremental math.",
    "Result: End-to-end detection latency held under 0.80 ms at scale."
)

# Slide 2
add_slide(
    "C O M P L I A N C E   &   I N N O V A T I O N   ·   P I L L A R   2   &   3",
    "Zero LLM hallucinations in compliance.",
    "Graceful degradation and deterministic NLG templates.",
    "The LLM Hallucination Risk",
    "Using public LLMs to write Suspicious Activity Reports (SARs) risks exposing PII (Data Privacy breach) and hallucinating facts (Legal liability).",
    "Banks cannot afford \"black box\" compliance documents that invent details about customers.",
    "Deterministic XAI Generation",
    "We engineered a Deterministic NLG Template Engine. It maps exact SHAP mathematical outputs directly to English sentences.",
    "Additionally, if the API for our conversational AI Copilot fails, the system gracefully degrades to a local \"Quick Commands\" intent router.",
    "Result: 100% factual accuracy, zero PII leakage, zero downtime."
)

# Slide 3
add_slide(
    "U S E R   E X P E R I E N C E   ·   P I L L A R   3",
    "The \"Single Pane of Glass\" workflow.",
    "Zero cognitive overload for investigators.",
    "The Cognitive Overload",
    "Today, compliance officers juggle four different screens: Core Banking Ledgers, KYC Portals, Network Analysis Tools, and Microsoft Word.",
    "Switching context causes errors and destroys productivity.",
    "Unified Single Pane of Glass",
    "We designed the React UX to hide the complex math from the end-user. We translate GraphSAGE metrics and SHAP values into simple English bullet points.",
    "With Live SSE streams and API View Caching, the UI refreshes instantly without loading spinners.",
    "Result: Detect, Explain, and File from one screen in real-time."
)

prs.save("c:/Users/satya/Desktop/bhadra_rudra/extra_slides.pptx")
print("Done")
